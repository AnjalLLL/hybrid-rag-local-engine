"""Build the local vector and keyword indexes from study files."""

from __future__ import annotations

import json
import pickle
import re
import site
import sys
from pathlib import Path
from typing import Dict, List, Sequence, Tuple

user_site = site.getusersitepackages()
if user_site and user_site not in sys.path:
    sys.path.append(user_site)

from src.utils import (
    DATA_DIR,
    FAISS_DIR,
    MIN_CHUNK_TOKENS,
    MODEL_NAME,
    chunk_text,
    cleanup,
    configure_offline_runtime,
    count_tokens,
    load_embedding_model,
)


ChunkRecord = Tuple[str, Dict[str, int | str]]
BATCH_SIZE = 64
SUPPORTED_EXTENSIONS = {".pdf", ".pptx", ".r", ".rmd", ".txt", ".md"}


def tokenize_text(text: str) -> List[str]:
    """Tokenize text for BM25 while preserving code-friendly terms."""

    return re.findall(r"[a-zA-Z_][a-zA-Z0-9_.]*|[0-9]+", text.lower())


def build_embed_text(text: str, source: str) -> str:
    """Prefix a chunk with its document title so retrieval can match on it.

    Slide decks reference their topic only in the filename, so a bare chunk
    body often has nothing for the query to match against.
    """

    title = Path(source).stem.replace("_", " ").replace("-", " ")
    return f"{title}\n{text}"


def discover_source_files(data_dir: Path = DATA_DIR) -> List[Path]:
    """Return supported study files found recursively under the data directory."""

    return sorted(
        path
        for path in data_dir.rglob("*")
        if path.is_file() and path.suffix.lower() in SUPPORTED_EXTENSIONS
    )


def extract_chunks_from_files(file_paths: Sequence[Path]) -> Tuple[List[ChunkRecord], List[str]]:
    """Extract, chunk, and annotate text from supported study files."""

    chunks: List[ChunkRecord] = []
    empty_sources: List[str] = []
    chunk_id = 0

    for file_path in file_paths:
        suffix = file_path.suffix.lower()
        if suffix == ".pdf":
            records = extract_chunks_from_pdf(file_path, chunk_id)
        elif suffix == ".pptx":
            records = extract_chunks_from_pptx(file_path, chunk_id)
        else:
            records = extract_chunks_from_text_file(file_path, chunk_id)

        if not records:
            empty_sources.append(file_path.name)

        chunks.extend(records)
        chunk_id += len(records)

    return chunks, empty_sources


def _emit_chunks(
    text: str,
    source: str,
    page: int,
    kind: str,
    chunk_id: int,
) -> List[ChunkRecord]:
    """Turn one page/slide/file of text into annotated chunk records."""

    records: List[ChunkRecord] = []
    for chunk in chunk_text(text):
        records.append(
            (
                chunk,
                {
                    "source": source,
                    "page": page,
                    "kind": kind,
                    "chunk_id": chunk_id + len(records),
                },
            )
        )
    return records


def extract_chunks_from_pdf(pdf_path: Path, start_chunk_id: int = 0) -> List[ChunkRecord]:
    """Extract chunk records from one PDF.

    Pages too small to chunk on their own are carried into the next page so their
    text still reaches the index, and the citation keeps the first page they came from.
    """

    import fitz

    chunks: List[ChunkRecord] = []
    chunk_id = start_chunk_id
    pending_text = ""
    pending_page = 1

    with fitz.open(pdf_path) as document:
        for page_index, page in enumerate(document, start=1):
            text = page.get_text("text").strip()
            if not text:
                continue

            if not pending_text:
                pending_page = page_index
            pending_text = f"{pending_text}\n{text}".strip() if pending_text else text

            if count_tokens(pending_text) < MIN_CHUNK_TOKENS:
                continue

            records = _emit_chunks(pending_text, pdf_path.name, pending_page, "pdf", chunk_id)
            chunks.extend(records)
            chunk_id += len(records)
            pending_text = ""

    if pending_text:
        records = _emit_chunks(pending_text, pdf_path.name, pending_page, "pdf", chunk_id)
        chunks.extend(records)

    return chunks


def extract_chunks_from_pptx(pptx_path: Path, start_chunk_id: int = 0) -> List[ChunkRecord]:
    """Extract chunk records from one PowerPoint file."""

    from pptx import Presentation

    chunks: List[ChunkRecord] = []
    chunk_id = start_chunk_id
    pending_text = ""
    pending_slide = 1
    presentation = Presentation(pptx_path)

    for slide_index, slide in enumerate(presentation.slides, start=1):
        text_parts = [shape.text for shape in slide.shapes if getattr(shape, "text", "").strip()]
        text = "\n".join(text_parts).strip()
        if not text:
            continue

        if not pending_text:
            pending_slide = slide_index
        pending_text = f"{pending_text}\n{text}".strip() if pending_text else text

        if count_tokens(pending_text) < MIN_CHUNK_TOKENS:
            continue

        records = _emit_chunks(pending_text, pptx_path.name, pending_slide, "pptx", chunk_id)
        chunks.extend(records)
        chunk_id += len(records)
        pending_text = ""

    if pending_text:
        records = _emit_chunks(pending_text, pptx_path.name, pending_slide, "pptx", chunk_id)
        chunks.extend(records)

    return chunks


def extract_chunks_from_text_file(text_path: Path, start_chunk_id: int = 0) -> List[ChunkRecord]:
    """Extract chunk records from one plain-text or R source file."""

    text = text_path.read_text(encoding="utf-8", errors="ignore").strip()
    if not text:
        return []

    return _emit_chunks(
        text,
        text_path.name,
        1,
        text_path.suffix.lower().lstrip("."),
        start_chunk_id,
    )


def embed_texts(model, texts: Sequence[str], batch_size: int = BATCH_SIZE):
    """Embed texts in batches and return normalized float32 vectors."""

    import numpy as np

    embeddings = model.encode(
        list(texts),
        batch_size=batch_size,
        convert_to_numpy=True,
        normalize_embeddings=True,
        show_progress_bar=True,
    )
    return np.asarray(embeddings, dtype="float32")


def persist_index(embeddings, chunks: Sequence[ChunkRecord], faiss_dir: Path = FAISS_DIR) -> None:
    """Write the FAISS index and chunk metadata to disk."""

    import faiss
    import numpy as np

    faiss_dir.mkdir(parents=True, exist_ok=True)
    index = faiss.IndexFlatIP(embeddings.shape[1])
    index.add(embeddings)
    faiss.write_index(index, str(faiss_dir / "index.faiss"))
    np.save(faiss_dir / "embeddings.npy", embeddings)

    payload = [{"text": text, "metadata": metadata} for text, metadata in chunks]
    tokenized_corpus = [
        tokenize_text(build_embed_text(str(chunk["text"]), str(chunk["metadata"]["source"])))
        for chunk in payload
    ]
    with (faiss_dir / "chunks.pkl").open("wb") as handle:
        pickle.dump({"chunks": payload, "tokenized_corpus": tokenized_corpus}, handle)

    meta = {
        "embedding_model": MODEL_NAME,
        "dimension": int(embeddings.shape[1]),
        "num_chunks": len(payload),
    }
    (faiss_dir / "meta.json").write_text(json.dumps(meta, indent=2), encoding="utf-8")


def build_index(rebuild: bool = False) -> int:
    """Build the FAISS index from the study files under the data directory."""

    configure_offline_runtime()
    index_path = FAISS_DIR / "index.faiss"
    chunks_path = FAISS_DIR / "chunks.pkl"

    if not rebuild and index_path.exists() and chunks_path.exists():
        print(f"Index already exists at {FAISS_DIR}. Use --rebuild to force a fresh build.")
        return 0

    source_paths = discover_source_files()
    if not source_paths:
        print("No supported files found under ./data/. Add files and run `python main.py ingest` again.")
        return 1

    print(f"Found {len(source_paths)} study files. Extracting text...")
    chunks, empty_sources = extract_chunks_from_files(source_paths)
    if not chunks:
        print("No usable text chunks were extracted from ./data/.")
        return 1

    if empty_sources:
        print(
            f"\nWarning: {len(empty_sources)} file(s) yielded no extractable text and are NOT in the index.\n"
            "These are most likely scanned/image-only PDFs that need OCR:"
        )
        for name in empty_sources:
            print(f"  - {name}")
        print()

    print(f"Building embeddings for {len(chunks)} chunks...")
    model = load_embedding_model()
    texts = [build_embed_text(text, str(metadata["source"])) for text, metadata in chunks]
    embeddings = embed_texts(model, texts)

    print("Writing FAISS index...")
    persist_index(embeddings, chunks)
    cleanup()
    print(f"Index build complete. Saved artifacts to {FAISS_DIR}.")
    return 0


if __name__ == "__main__":
    configure_offline_runtime()
    sys.exit(build_index(rebuild=False))
